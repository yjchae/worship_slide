import concurrent.futures
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

PPT_CONVERT_CHUNK_SIZE = 40
PPT_CONVERT_PARALLELISM = 3

_PRETENDARD_FONTS = ["Pretendard-Bold.ttf", "Pretendard-Regular.ttf"]


def _get_font_path(filename):
    if getattr(sys, "frozen", False):
        return os.path.join(sys._MEIPASS, "fonts", filename)
    script_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(script_dir, "..", "assets", "fonts", filename)


def _ensure_pretendard_installed():
    import platform
    system = platform.system()
    if system == "Windows":
        local_appdata = os.environ.get("LOCALAPPDATA", "")
        if not local_appdata:
            return
        fonts_dir = os.path.join(local_appdata, "Microsoft", "Windows", "Fonts")
    else:
        fonts_dir = os.path.expanduser("~/Library/Fonts")
    os.makedirs(fonts_dir, exist_ok=True)
    for filename in _PRETENDARD_FONTS:
        dest = os.path.join(fonts_dir, filename)
        if not os.path.exists(dest):
            src = _get_font_path(filename)
            if os.path.exists(src):
                shutil.copy2(src, dest)


def extract_text_from_shape(shape):
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return ""
    lines = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            lines.append(text)
    return "\n".join(lines)


def _find_title_texts(prs):
    """텍스트가 있는 슬라이드의 80% 이상에서 동일한 전체 텍스트로 반복되는
    텍스트박스의 내용을 반환한다. (위치가 아닌 shape 전체 텍스트의 완전 일치 기준)
    가사 shape는 슬라이드마다 다른 내용을 가지므로 걸러지지 않고,
    제목 shape만 동일 텍스트로 매 슬라이드에 반복되어 걸러진다."""
    text_count = {}
    content_slide_count = 0

    for slide in prs.slides:
        seen = set()
        slide_has_text = False
        for shape in slide.shapes:
            text = extract_text_from_shape(shape).strip()
            if not text:
                continue
            slide_has_text = True
            if text not in seen:
                seen.add(text)
                text_count[text] = text_count.get(text, 0) + 1
        if slide_has_text:
            content_slide_count += 1

    if content_slide_count < 2:
        return set()

    threshold = max(2, round(content_slide_count * 0.8))
    return {text for text, count in text_count.items() if count >= threshold}


def _shape_max_font_pt(shape):
    """명시적으로 설정된 최대 폰트 크기(pt)를 반환. 없으면 None."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return None
    max_pt = None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                pt = run.font.size.pt
                max_pt = pt if max_pt is None else max(max_pt, pt)
    return max_pt


def slide_lyrics(slide, title_texts=None):
    shapes_info = []
    for shape in slide.shapes:
        text = extract_text_from_shape(shape)
        if not text:
            continue
        # shape 전체 텍스트가 반복 제목과 완전히 일치하는 shape만 제외
        if title_texts and text.strip() in title_texts:
            continue
        shapes_info.append((text, _shape_max_font_pt(shape)))

    if not shapes_info:
        return ""
    if len(shapes_info) == 1:
        return shapes_info[0][0]

    # 보조 필터: 명시적 폰트 크기가 최댓값의 60% 미만인 박스도 제목/캡션으로 간주
    explicit_sizes = [size for _, size in shapes_info if size is not None]
    if explicit_sizes:
        max_size = max(explicit_sizes)
        kept = [
            text for text, size in shapes_info
            if size is None or size >= max_size * 0.6
        ]
        if kept:
            return "\n".join(kept).strip()

    return "\n".join(text for text, _ in shapes_info).strip()


def normalize_title(path):
    # macOS HFS+는 파일명을 NFD로 저장하므로 NFC로 변환해야 한글이 정상 표시됨
    return unicodedata.normalize('NFC', path.stem.strip())


def is_english_line(line):
    stripped = line.strip()
    if not stripped:
        return False
    if any("\uac00" <= char <= "\ud7a3" for char in stripped):
        return False

    alpha_count = sum(char.isalpha() for char in stripped)
    latin_count = sum(("A" <= char <= "Z") or ("a" <= char <= "z") for char in stripped)
    if latin_count == 0:
        return False
    return latin_count / max(alpha_count, 1) >= 0.6


def split_bilingual_page(text):
    korean_lines = []
    english_lines = []

    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if is_english_line(stripped):
            english_lines.append(stripped)
        else:
            korean_lines.append(stripped)

    return "\n".join(korean_lines), "\n".join(english_lines)


def get_cache_root():
    import platform
    if platform.system() == "Windows":
        local_appdata = os.environ.get("LOCALAPPDATA")
        if local_appdata:
            return Path(local_appdata) / "worship_slides" / "ppt_import_cache"
        return Path.home() / "AppData" / "Local" / "worship_slides" / "ppt_import_cache"
    return Path.home() / "Library" / "Caches" / "worship_slides" / "ppt_import_cache"


def get_ppt_cache_key(source_path):
    stat = source_path.stat()
    payload = f"{source_path.resolve()}::{stat.st_size}::{stat.st_mtime_ns}"
    return hashlib.sha1(payload.encode("utf-8")).hexdigest()


def get_cached_pptx_path(source_path):
    cache_root = get_cache_root()
    cache_root.mkdir(parents=True, exist_ok=True)
    return cache_root / f"{get_ppt_cache_key(source_path)}.pptx"


def get_libreoffice_executable():
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/LibreOffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if candidate and os.path.exists(candidate):
            return candidate
    return None


def convert_ppt_group_to_pptx(source_paths, output_dir):
    libreoffice = get_libreoffice_executable()
    if libreoffice is None:
        raise RuntimeError("LibreOffice 실행 파일을 찾지 못했습니다.")

    subprocess.run(
        [
            libreoffice,
            "--headless",
            "--convert-to",
            "pptx",
            "--outdir",
            str(output_dir),
            *[str(source_path) for source_path in source_paths],
        ],
        check=True,
        capture_output=True,
        text=True,
    )

    converted_map = {}
    used_targets = set()
    for source_path in source_paths:
        stem_pattern = re.escape(source_path.stem)
        candidates = sorted(output_dir.glob(f"{source_path.stem}*.pptx"))
        matched = next(
            (
                candidate
                for candidate in candidates
                if candidate not in used_targets
                and re.fullmatch(rf"{stem_pattern}(?:_[0-9]+)?\.pptx", candidate.name, re.IGNORECASE)
            ),
            None,
        )
        if matched is not None:
            used_targets.add(matched)
            converted_map[source_path] = matched

    return converted_map


def chunked(items, chunk_size):
    for index in range(0, len(items), chunk_size):
        yield items[index : index + chunk_size]


def iter_presentation_files(root):
    def handle_walk_error(error):
        print(
            json.dumps(
                {
                    "warning": f"폴더를 읽지 못했습니다: {error.filename}",
                },
                ensure_ascii=True,
            ),
            file=sys.stderr,
        )

    for current_root, dir_names, file_names in os.walk(root, onerror=handle_walk_error):
        dir_names[:] = [name for name in dir_names if not name.startswith(".")]
        for file_name in sorted(file_names):
            suffix = Path(file_name).suffix.lower()
            if suffix not in {".ppt", ".pptx"}:
                continue
            yield Path(current_root) / file_name


def process_presentation_file(file_path, presentation_path):
    try:
        prs = Presentation(str(presentation_path))
        title_texts = _find_title_texts(prs)
        korean_pages = []
        english_pages = []

        for slide in prs.slides:
            page = slide_lyrics(slide, title_texts=title_texts)
            if not page:
                continue
            korean_page, english_page = split_bilingual_page(page)
            if korean_page:
                korean_pages.append(korean_page)
            else:
                korean_pages.append("")
            english_pages.append(english_page)

        if not any(page for page in korean_pages + english_pages):
            return {"status": "skipped", "path": str(file_path)}

        return {
            "status": "ok",
            "song": {
                "file_name": file_path.name,
                "title": normalize_title(file_path),
                "lyrics": "\n###\n".join(korean_pages),
                "english_lyrics": "\n###\n".join(english_pages),
            },
            "path": str(file_path),
        }
    except Exception as error:
        return {
            "status": "error",
            "error": {
                "file_name": file_path.name,
                "path": str(file_path),
                "error": str(error),
            },
        }


def _is_libreoffice_available():
    return get_libreoffice_executable() is not None


def prepare_presentation_sources(file_paths):
    prepared_sources = []
    errors = []
    temp_dirs = []
    ppt_groups = {}

    for file_path in file_paths:
        if file_path.suffix.lower() == ".ppt":
            cached_pptx = get_cached_pptx_path(file_path)
            if cached_pptx.exists():
                prepared_sources.append((file_path, cached_pptx))
            else:
                ppt_groups.setdefault(file_path.parent, []).append(file_path)
        else:
            prepared_sources.append((file_path, file_path))

    if ppt_groups and not _is_libreoffice_available():
        for file_path_list in ppt_groups.values():
            for file_path in file_path_list:
                errors.append(
                    {
                        "file_name": file_path.name,
                        "path": str(file_path),
                        "error": "libreoffice_missing",
                    }
                )
        ppt_groups = {}

    for index, directory in enumerate(sorted(ppt_groups, key=lambda item: str(item).casefold())):
        source_paths = sorted(ppt_groups[directory], key=lambda path: str(path).casefold())
        chunks = list(chunked(source_paths, PPT_CONVERT_CHUNK_SIZE))
        chunk_jobs = []
        for chunk_index, source_chunk in enumerate(chunks):
            output_dir = Path(
                tempfile.mkdtemp(prefix=f"praise_ppt_batch_{index}_{chunk_index}_")
            )
            temp_dirs.append(output_dir)
            chunk_jobs.append((source_chunk, output_dir))

        with concurrent.futures.ThreadPoolExecutor(
            max_workers=min(PPT_CONVERT_PARALLELISM, len(chunk_jobs))
        ) as executor:
            future_map = {
                executor.submit(convert_ppt_group_to_pptx, source_chunk, output_dir): (
                    source_chunk,
                    output_dir,
                )
                for source_chunk, output_dir in chunk_jobs
            }
            for future in concurrent.futures.as_completed(future_map):
                source_chunk, _ = future_map[future]
                try:
                    converted_map = future.result()
                    for source_path in source_chunk:
                        converted_path = converted_map.get(source_path)
                        if converted_path is None or not converted_path.exists():
                            errors.append(
                                {
                                    "file_name": source_path.name,
                                    "path": str(source_path),
                                    "error": f"PPT 변환 실패: {source_path}",
                                }
                            )
                            continue

                        cached_pptx = get_cached_pptx_path(source_path)
                        shutil.copy2(converted_path, cached_pptx)
                        prepared_sources.append((source_path, cached_pptx))
                except Exception as error:
                    for source_path in source_chunk:
                        errors.append(
                            {
                                "file_name": source_path.name,
                                "path": str(source_path),
                                "error": str(error),
                            }
                        )

    prepared_sources.sort(key=lambda item: str(item[0]).casefold())
    return prepared_sources, errors, temp_dirs


def get_import_worker_count(file_count):
    cpu_count = os.cpu_count() or 4
    return max(1, min(file_count, cpu_count, 8))


def import_folder(folder):
    root = Path(folder)
    if not root.exists():
        raise RuntimeError(f"폴더를 찾을 수 없습니다: {folder}")

    file_paths = list(iter_presentation_files(root))
    songs = []
    errors = []
    processed_count = len(file_paths)

    prepared_sources, preparation_errors, temp_dirs = prepare_presentation_sources(file_paths)
    errors.extend(preparation_errors)

    try:
        if prepared_sources:
            max_workers = get_import_worker_count(len(prepared_sources))
            with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = [
                    executor.submit(process_presentation_file, source_path, presentation_path)
                    for source_path, presentation_path in prepared_sources
                ]
                for future in concurrent.futures.as_completed(futures):
                    result = future.result()
                    if result["status"] == "ok":
                        songs.append(result["song"])
                    elif result["status"] == "error":
                        errors.append(result["error"])
    finally:
        for temp_dir in temp_dirs:
            shutil.rmtree(temp_dir, ignore_errors=True)

    songs.sort(key=lambda song: (song["title"].casefold(), song["file_name"].casefold()))
    errors.sort(key=lambda entry: entry["path"].casefold())

    print(
        json.dumps(
            {
                "songs": songs,
                "processed_count": processed_count,
                "imported_count": len(songs),
                "failed_count": len(errors),
                "errors": errors,
                "worker_count": get_import_worker_count(processed_count),
                "libreoffice_missing": any(
                    e["error"] == "libreoffice_missing" for e in errors
                ),
            },
            ensure_ascii=True,
        )
    )


def parse_hex_color(hex_color):
    value = hex_color.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


_TITLE_BOX_HEIGHT = 0.55
_TITLE_BOX_PADDING = 0.2
_TITLE_BOX_WIDTH_CENTER = 10.0
_SLIDE_W = 13.333
_SLIDE_H = 7.5
_TITLE_BOX_WIDTH_SIDE = _SLIDE_W - (_TITLE_BOX_PADDING * 2)
_LYRICS_BOX_TOP = 0.6
_LYRICS_BOX_HEIGHT = 5.4
_LYRICS_BOX_BOTTOM = _SLIDE_H - _LYRICS_BOX_TOP - _LYRICS_BOX_HEIGHT
_LYRICS_BOX_WIDTH = _SLIDE_W * 0.9
_LYRICS_BOX_LEFT = (_SLIDE_W - _LYRICS_BOX_WIDTH) / 2
_TEXT_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}
_FONT_NAME = "Pretendard"


def _lyrics_text_layout(horizontal_position):
    text_align = _TEXT_ALIGN_MAP.get(horizontal_position, PP_ALIGN.CENTER)
    return _LYRICS_BOX_LEFT, _LYRICS_BOX_WIDTH, text_align


def _lyrics_box_vertical_layout(style, is_bible):
    top = float(style.get(
        "bible_text_box_top" if is_bible else "text_box_top",
        _LYRICS_BOX_TOP,
    ))
    top = min(max(top, 0.3), _SLIDE_H - _LYRICS_BOX_BOTTOM - 1.0)
    return top, _SLIDE_H - top - _LYRICS_BOX_BOTTOM


def _add_title_textbox(slide, song_title, style, is_bible=False):
    title_font_size = Pt(style.get(
        "bible_title_font_size" if is_bible else "title_font_size",
        style.get("title_font_size", 14),
    ))
    title_color = parse_hex_color(style.get(
        "bible_title_text_color" if is_bible else "title_text_color",
        style.get("title_text_color", "#B3FFFFFF"),
    ))
    h_pos = style.get(
        "bible_title_horizontal_position" if is_bible else "title_horizontal_position",
        style.get("title_horizontal_position", "right"),
    )
    v_pos = style.get(
        "bible_title_vertical_position" if is_bible else "title_vertical_position",
        style.get("title_vertical_position", "bottom"),
    )

    if h_pos == "left":
        title_left = _TITLE_BOX_PADDING
        title_width = _TITLE_BOX_WIDTH_SIDE
        text_align = PP_ALIGN.LEFT
    elif h_pos == "center":
        title_left = (_SLIDE_W - _TITLE_BOX_WIDTH_CENTER) / 2
        title_width = _TITLE_BOX_WIDTH_CENTER
        text_align = PP_ALIGN.CENTER
    else:
        title_left = _SLIDE_W - _TITLE_BOX_PADDING - _TITLE_BOX_WIDTH_SIDE
        title_width = _TITLE_BOX_WIDTH_SIDE
        text_align = PP_ALIGN.RIGHT

    if v_pos == "top":
        title_top = _TITLE_BOX_PADDING
    elif v_pos == "middle":
        title_top = (_SLIDE_H - _TITLE_BOX_HEIGHT) / 2
    else:
        title_top = _SLIDE_H - _TITLE_BOX_PADDING - _TITLE_BOX_HEIGHT

    box = slide.shapes.add_textbox(
        Inches(title_left), Inches(title_top),
        Inches(title_width), Inches(_TITLE_BOX_HEIGHT),
    )
    frame = box.text_frame
    frame.word_wrap = False
    para = frame.paragraphs[0]
    para.alignment = text_align
    run = para.add_run()
    run.text = _normalize_ppt_text(song_title)
    _set_run_font(run)
    run.font.size = title_font_size
    run.font.color.rgb = title_color


def _normalize_ppt_text(text):
    return unicodedata.normalize("NFC", text or "")


def _set_run_font(run):
    run.font.name = _FONT_NAME
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        font = rpr.find(qn(tag))
        if font is None:
            font = OxmlElement(tag)
            rpr.append(font)
        font.set("typeface", _FONT_NAME)


def _add_text_run(paragraph, text, font_size, color):
    run = paragraph.add_run()
    run.text = _normalize_ppt_text(text)
    _set_run_font(run)
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = color


def _format_bible_paragraph(paragraph, text_align, font_size):
    paragraph.alignment = text_align
    hanging_width = font_size.pt * 1.7
    paragraph.margin_left = Pt(hanging_width)
    paragraph.first_line_indent = Pt(-hanging_width)


def _add_bible_page_text(frame, page, text_align, font_size, text_color):
    lines = [line.strip() for line in page.splitlines() if line.strip()]
    if not lines:
        return

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        _format_bible_paragraph(paragraph, text_align, font_size)
        _add_text_run(paragraph, line, font_size, text_color)


def add_song_slides(prs, song, style):
    is_bible = song.get("type") == "bible"
    bg_color = parse_hex_color(style["background_color"])
    text_color_key = "bible_text_color" if is_bible else "text_color"
    text_color = parse_hex_color(style.get(text_color_key, style["text_color"]))
    font_size = Pt(style.get(
        "bible_font_size" if is_bible else "font_size",
        style["font_size"],
    ))
    position = style.get(
        "bible_text_position" if is_bible else "text_position",
        style["text_position"],
    )
    include_english_lyrics = style.get("include_english_lyrics", False)
    english_color = parse_hex_color(style["english_text_color"])
    show_song_title = style.get(
        "show_bible_title" if is_bible else "show_song_title",
        style.get("show_song_title", False),
    )
    lyrics_text_align_str = style.get(
        "bible_text_align" if is_bible else "lyrics_text_align", "center"
    )
    lyrics_box_left, lyrics_box_width, lyrics_align = _lyrics_text_layout(
        lyrics_text_align_str
    )
    korean_pages = [part.strip() for part in song["lyrics"].split("###")]
    english_pages = [part.strip() for part in song.get("english_lyrics", "").split("###")]
    page_count = max(len(korean_pages), len(english_pages))

    for index in range(page_count):
        page = korean_pages[index] if index < len(korean_pages) else ""
        english_page = english_pages[index] if index < len(english_pages) else ""
        if not page and not english_page:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        lyrics_box_top, lyrics_box_height = _lyrics_box_vertical_layout(
            style, is_bible
        )

        textbox = slide.shapes.add_textbox(
            Inches(lyrics_box_left), Inches(lyrics_box_top),
            Inches(lyrics_box_width), Inches(lyrics_box_height),
        )
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[position]

        if is_bible:
            _add_bible_page_text(frame, page, lyrics_align, font_size, text_color)
        else:
            paragraph = frame.paragraphs[0]
            paragraph.alignment = lyrics_align
            _add_text_run(paragraph, page, font_size, text_color)

        if include_english_lyrics and english_page:
            english_paragraph = frame.add_paragraph()
            english_paragraph.alignment = lyrics_align
            _add_text_run(
                english_paragraph,
                english_page,
                Pt(style["font_size"] * 0.8),
                english_color,
            )

        if show_song_title:
            _add_title_textbox(slide, song.get("title", ""), style, is_bible=is_bible)


def _add_blank_slide(prs, style):
    bg_color = parse_hex_color(style["background_color"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color


def export_presentation(payload_json):
    _ensure_pretendard_installed()
    payload = json.loads(payload_json)
    output_path = Path(payload["output_path"])
    if output_path.suffix.lower() != ".pptx":
        output_path = output_path.with_suffix(".pptx")
    songs = payload["songs"]
    style = payload["style"]

    prs = Presentation()
    if len(prs.slides) == 0:
        prs.slides.add_slide(prs.slide_layouts[6])
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, song in enumerate(songs):
        add_song_slides(prs, song, style)
        if index < len(songs) - 1:
            _add_blank_slide(prs, style)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(json.dumps({"output_path": str(output_path)}, ensure_ascii=True))


def main():
    if len(sys.argv) < 3:
        raise SystemExit("usage: ppt_tool.py [import|export] [payload]")

    command = sys.argv[1]
    payload = sys.argv[2]

    if command == "import":
        import_folder(payload)
        return

    if command == "export":
        export_presentation(payload)
        return

    raise SystemExit(f"unknown command: {command}")


if __name__ == "__main__":
    main()
