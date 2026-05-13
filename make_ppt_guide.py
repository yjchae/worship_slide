"""예배 슬라이드 사용자 가이드 PPT 생성."""

import io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ── 폰트 경로 ────────────────────────────────────────────────────────────────
_FONT_DIR = '/home/user/make_ppt/assets/fonts'
_BOLD = f'{_FONT_DIR}/Pretendard-Bold.ttf'
_REG  = f'{_FONT_DIR}/Pretendard-Regular.ttf'

def pf(size, bold=False):
    return ImageFont.truetype(_BOLD if bold else _REG, size)

# ── 슬라이드 크기: 와이드 16:9 (33.87 x 19.05 cm) ──────────────────────────
SW = Cm(33.87)
SH = Cm(19.05)

# ── 색상 팔레트 ──────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x14, 0x36, 0x42)
C_TEAL   = RGBColor(0x0F, 0x8B, 0x8D)
C_LTEAL  = RGBColor(0xC8, 0xEC, 0xEC)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x60, 0x60, 0x65)
C_LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x0F, 0x4C, 0x5C)
C_YELLOW = RGBColor(0xFF, 0xF1, 0x76)

# PIL 튜플
PT_DARK  = (0x14, 0x36, 0x42)
PT_TEAL  = (0x0F, 0x8B, 0x8D)
PT_LTEAL = (0xC8, 0xEC, 0xEC)
PT_WHITE = (0xFF, 0xFF, 0xFF)
PT_LGRAY = (0xF4, 0xF6, 0xF8)
PT_GRAY  = (0x90, 0x90, 0x96)
PT_BLUE  = (0x18, 0x7A, 0xB8)

# ── PPT 헬퍼 ─────────────────────────────────────────────────────────────────

def new_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def textbox(slide, text, x, y, w, h,
            font_size=18, bold=False, color: RGBColor = C_WHITE,
            align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Pretendard'
    return txb


def add_image(slide, img: Image.Image, x, y, w, h=None):
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    if h:
        slide.shapes.add_picture(buf, x, y, w, h)
    else:
        slide.shapes.add_picture(buf, x, y, w)


def pill(slide, text, x, y, w, h, bg_color: RGBColor,
         text_color: RGBColor = C_WHITE, font_size=14, bold=False):
    r = rect(slide, x, y, w, h, bg_color)
    r.line.fill.background()
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = 'Pretendard'
    return r

# ── PIL 이미지 유틸 ───────────────────────────────────────────────────────────

def gradient_rect(d, x0, y0, x1, y1, c1, c2, vertical=True):
    steps = (y1 - y0) if vertical else (x1 - x0)
    for i in range(steps):
        t = i / max(steps - 1, 1)
        r = int(c1[0] + (c2[0] - c1[0]) * t)
        g = int(c1[1] + (c2[1] - c1[1]) * t)
        b = int(c1[2] + (c2[2] - c1[2]) * t)
        if vertical:
            d.line([(x0, y0 + i), (x1, y0 + i)], fill=(r, g, b))
        else:
            d.line([(x0 + i, y0), (x0 + i, y1)], fill=(r, g, b))


# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드별 이미지 생성
# ══════════════════════════════════════════════════════════════════════════════

def img_app_overview() -> Image.Image:
    """앱 기능 요약 아이콘 카드."""
    W, H = 1200, 420
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    cards = [
        ("📂", "찬양 폴더\n가져오기", "PPT/PPTX 파일을\n한번에 등록"),
        ("🔍", "찬양 검색\n& 선택",   "곡명·가사로\n빠르게 찾기"),
        ("📖", "성경 본문\n추가",     "책·장·절 선택\n슬라이드 생성"),
        ("🎨", "디자인\n설정",        "색상·크기·위치\n미리보기"),
        ("📽", "화면 발표",           "외부 스크린\n전체화면 송출"),
    ]
    card_w, card_h = 210, 340
    gap = 18
    total = len(cards) * card_w + (len(cards) - 1) * gap
    sx = (W - total) // 2

    for i, (icon, title, sub) in enumerate(cards):
        x = sx + i * (card_w + gap)
        # 카드 배경
        d.rounded_rectangle([x, 30, x + card_w, x + card_h + 30 - x + x],
                              radius=16, fill=(255, 255, 255),
                              outline=(220, 225, 230))
        # 아이콘 원
        cx = x + card_w // 2
        d.ellipse([cx - 38, 60, cx + 38, 136],
                   fill=PT_TEAL if i % 2 == 0 else PT_DARK)
        # 아이콘 텍스트는 넣기 어려우므로 번호로 대체 (색 원)
        num_font = pf(22, bold=True)
        d.text((cx - 8, 82), str(i + 1), fill=PT_WHITE, font=num_font)
        # 제목
        tf = pf(17, bold=True)
        for j, line in enumerate(title.split('\n')):
            d.text((x + 20, 152 + j * 26), line, fill=PT_DARK, font=tf)
        # 설명
        sf = pf(13)
        for j, line in enumerate(sub.split('\n')):
            d.text((x + 20, 212 + j * 20), line, fill=(100, 110, 115), font=sf)

    return img


def img_main_screen() -> Image.Image:
    """메인 화면 레이아웃 설명."""
    W, H = 1400, 820
    img = Image.new('RGB', (W, H), (238, 242, 246))
    d = ImageDraw.Draw(img)

    # 앱 프레임
    d.rounded_rectangle([20, 20, W - 20, H - 20], radius=14,
                          fill=(250, 250, 252), outline=(200, 205, 212))

    # 상단 헤더
    for i in range(62):
        t = i / 62
        r = int(PT_DARK[0] + ((0x0F - PT_DARK[0]) * t))
        g = int(PT_DARK[1] + ((0x8B - PT_DARK[1]) * t))
        b = int(PT_DARK[2] + ((0x8D - PT_DARK[2]) * t))
        d.line([(22, 22 + i), (W - 22, 22 + i)], fill=(r, g, b))
    d.rounded_rectangle([22, 22, W - 22, 84], radius=14, outline=(0x0F, 0x8B, 0x8D))

    d.text((40, 32), "예배 슬라이드 보관함", fill=PT_WHITE, font=pf(18, True))
    d.text((290, 36), "저장 42곡", fill=(200, 240, 240), font=pf(13))
    d.rounded_rectangle([W - 340, 32, W - 180, 72], radius=6,
                          fill=(255, 255, 255, 50))
    d.text((W - 328, 40), "찬양폴더 선택", fill=PT_WHITE, font=pf(14, True))
    d.rounded_rectangle([W - 168, 32, W - 30, 72], radius=6,
                          fill=(255, 255, 255, 50))
    d.text((W - 156, 40), "성경 불러오기", fill=PT_WHITE, font=pf(14, True))

    # 발표 컨트롤 바
    d.rounded_rectangle([30, 92, W - 30, 140], radius=8, fill=(228, 234, 240))
    d.text((48, 106), "발표 화면", fill=(70, 80, 90), font=pf(14))
    d.rounded_rectangle([W - 200, 100, W - 38, 132], radius=6, fill=(60, 115, 170))
    d.text((W - 184, 107), "▶  발표 시작", fill=PT_WHITE, font=pf(14, True))

    # ── 왼쪽: 선택한 순서 패널 ──
    d.rounded_rectangle([30, 148, 760, 440], radius=10, fill=(255, 255, 255),
                          outline=(215, 218, 224))
    d.text((48, 160), "① 선택한 순서", fill=PT_DARK, font=pf(15, True))
    d.line([(48, 184), (742, 184)], fill=(220, 222, 228))
    rows = [
        ("1", "주님 한 분만으로", "주님 한 분만으로 내게는 족해", False),
        ("2", "주의 이름 높이세", "온 땅이여 주께 소리 질러", False),
        ("3", "로마서 8:28", "우리가 알거니와…", True),
    ]
    for i, (num, t, s, bi) in enumerate(rows):
        y = 192 + i * 76
        if i == 0:
            d.rectangle([31, y, 759, y + 72], fill=(237, 246, 255))
        d.text((50, y + 20), num, fill=(60, 65, 70), font=pf(14))
        if bi:
            d.rounded_rectangle([76, y + 16, 118, y + 38], radius=4, fill=(187, 222, 251))
            d.text((80, y + 20), "성경", fill=(21, 100, 190), font=pf(11))
            d.text((126, y + 19), t, fill=(30, 35, 40), font=pf(14))
        else:
            d.text((76, y + 19), t, fill=(30, 35, 40), font=pf(14))
        d.text((76, y + 42), s, fill=(148, 152, 158), font=pf(11))
        d.line([(32, y + 72), (758, y + 72)], fill=(228, 230, 236))

    # ── 왼쪽: 검색 패널 ──
    d.rounded_rectangle([30, 448, 760, H - 30], radius=10, fill=(255, 255, 255),
                          outline=(215, 218, 224))
    d.text((48, 462), "② 찬양 검색", fill=PT_DARK, font=pf(15, True))
    d.text((196, 462), "성경 본문", fill=(160, 162, 168), font=pf(14))
    d.line([(48, 484), (150, 484)], fill=PT_TEAL, width=2)
    d.rounded_rectangle([38, 492, 750, 528], radius=6, outline=(185, 188, 194),
                          fill=(252, 252, 254))
    d.text((52, 502), "곡명 또는 가사로 검색…", fill=(185, 188, 195), font=pf(13))
    songs = [("나는 예배자입니다", False), ("주님 한 분만으로", True),
             ("주의 이름 높이세", True), ("찬양하라 내 영혼아", False)]
    for i, (s, sel) in enumerate(songs):
        y = 534 + i * 52
        d.rectangle([31, y, 759, y + 48],
                      fill=(230, 248, 248) if sel else (255, 255, 255))
        d.text((52, y + 14),
               ("✓  " if sel else "    ") + s,
               fill=PT_TEAL if sel else (48, 52, 58), font=pf(13))
        d.line([(32, y + 48), (758, y + 48)], fill=(235, 237, 242))

    # ── 오른쪽: 디자인 패널 ──
    d.rounded_rectangle([776, 148, W - 30, H - 30], radius=10, fill=(255, 255, 255),
                          outline=(215, 218, 224))
    d.text((792, 160), "③ PPTX 디자인", fill=PT_DARK, font=pf(15, True))
    d.line([(792, 184), (W - 38, 184)], fill=(220, 222, 228))
    # 미리보기
    d.rounded_rectangle([784, 192, W - 38, 430], radius=6, fill=(27, 27, 27))
    d.text((860, 278), "주님 한 분만으로", fill=PT_WHITE, font=pf(22, True))
    d.text((860, 316), "Only You, Lord", fill=(255, 241, 118), font=pf(16))
    d.text((792, 436), "▲ 실시간 미리보기", fill=(128, 132, 140), font=pf(11))
    # 옵션
    opts = [("배경색", "■ 어두운"), ("글자색", "■ 흰색"), ("글자 크기", "30 pt"),
            ("영어 가사", "켜짐"), ("글자 위치", "중단")]
    for i, (lbl, val) in enumerate(opts):
        y = 454 + i * 34
        d.text((792, y), lbl, fill=(95, 98, 106), font=pf(13))
        d.text((930, y), val, fill=(30, 33, 40), font=pf(13))
    d.rounded_rectangle([784, 628, W - 38, 670], radius=8, fill=PT_TEAL)
    d.text((920, 638), "PPTX 저장", fill=PT_WHITE, font=pf(16, True))

    # ── 화살표 레이블 ──
    labels = [
        (390, 10, "상단 바: 폴더 가져오기 / 성경 불러오기"),
        (390, H - 14, "검색·성경 탭: 곡/성경 선택"),
        (1060, 10, "디자인 패널: 스타일 & 저장"),
    ]
    for lx, ly, lt in labels:
        d.text((lx - 100, ly), lt, fill=PT_TEAL, font=pf(12, True))

    return img


def img_import_steps() -> Image.Image:
    """찬양 가져오기 단계."""
    W, H = 1200, 500
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    steps = [
        ("1", "찬양폴더 선택\n버튼 클릭", PT_DARK),
        ("2", "PPT 폴더 선택\n(하위 폴더 포함)", (0x0F, 0x4C, 0x5C)),
        ("3", "자동 분석 중\n(수 초~수 분)", PT_TEAL),
        ("4", "저장 완료!\n검색 가능", (0x12, 0x96, 0x74)),
    ]
    bw, bh = 240, 280
    gap = 30
    total = len(steps) * bw + (len(steps) - 1) * gap
    sx = (W - total) // 2

    for i, (num, text, col) in enumerate(steps):
        x = sx + i * (bw + gap)
        gradient_rect(d, x, 60, x + bw, 60 + bh,
                       col, tuple(min(c + 30, 255) for c in col), vertical=True)
        d.rounded_rectangle([x, 60, x + bw, 60 + bh], radius=16, outline=(200, 210, 215))
        # 번호 원
        d.ellipse([x + bw // 2 - 28, 80, x + bw // 2 + 28, 136],
                   fill=(255, 255, 255, 80))
        d.text((x + bw // 2 - 10, 90), num,
               fill=PT_WHITE, font=pf(30, True))
        # 텍스트
        for j, line in enumerate(text.split('\n')):
            tw = len(line) * 14
            d.text((x + (bw - tw) // 2, 160 + j * 32), line,
                   fill=PT_WHITE, font=pf(17, True))
        # 화살표
        if i < len(steps) - 1:
            ax = x + bw + gap // 2 - 6
            ay = 60 + bh // 2
            d.polygon([(ax, ay - 14), (ax + 20, ay), (ax, ay + 14)],
                       fill=(0x0F, 0x8B, 0x8D))

    # 하단 팁
    tip_y = 380
    d.rounded_rectangle([sx, tip_y, sx + total, tip_y + 80], radius=10,
                          fill=(232, 248, 248), outline=(0x0F, 0x8B, 0x8D))
    d.text((sx + 20, tip_y + 22),
           "💡  같은 폴더를 여러 번 가져와도 안전해요 — 중복 곡은 자동으로 건너뜁니다",
           fill=PT_DARK, font=pf(14))

    return img


def img_search_select() -> Image.Image:
    """찬양 검색 & 선택 화면."""
    W, H = 1100, 560
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    # 검색 패널
    px, py, pw, ph = 40, 40, 500, 480
    d.rounded_rectangle([px, py, px + pw, py + ph], radius=12, fill=(255, 255, 255),
                          outline=(215, 218, 224))
    d.text((px + 16, py + 14), "찬양 검색", fill=PT_DARK, font=pf(16, True))
    d.line([(px + 16, py + 40), (px + 130, py + 40)], fill=PT_TEAL, width=2)
    # 검색창
    d.rounded_rectangle([px + 10, py + 50, px + pw - 10, py + 86], radius=6,
                          outline=(185, 188, 194), fill=(252, 252, 254))
    d.text((px + 24, py + 60), "주님", fill=PT_DARK, font=pf(14))
    # 목록
    songs = [
        ("주님 한 분만으로", True, "주님 한 분만으로 내게는 족해"),
        ("주님 다시 오실 때", False, "주님 다시 오실 때까지"),
        ("주님의 사랑이", False, "주님의 사랑이 이 땅에"),
    ]
    for i, (title, sel, sub) in enumerate(songs):
        y = py + 94 + i * 110
        d.rectangle([px + 1, y, px + pw - 1, y + 104],
                      fill=(232, 247, 247) if sel else (255, 255, 255))
        # 체크박스
        if sel:
            d.ellipse([px + 14, y + 14, px + 42, y + 42], fill=PT_TEAL)
            d.text((px + 19, y + 16), "✓", fill=PT_WHITE, font=pf(16, True))
        else:
            d.rounded_rectangle([px + 14, y + 14, px + 42, y + 42], radius=4,
                                  outline=(200, 205, 210))
        d.text((px + 56, y + 18), title,
               fill=PT_TEAL if sel else (30, 35, 40), font=pf(15, True if sel else False))
        d.text((px + 56, y + 46), sub, fill=(148, 152, 158), font=pf(12))
        d.text((px + 56, y + 68), f"{'→ 선택됨 (순서에 추가)' if sel else '클릭하면 체크'  }",
               fill=PT_TEAL if sel else (190, 195, 200), font=pf(11))
        d.line([(px + 2, y + 104), (px + pw - 2, y + 104)], fill=(232, 235, 240))

    # 오른쪽 설명
    rx = 580
    tips = [
        ("✓ 체크 = 선택한 순서에 추가", PT_TEAL),
        ("✓ 다시 클릭 = 해제 (목록 제거)", PT_TEAL),
        ("", None),
        ("검색창에 곡명이나 가사 일부를", PT_DARK),
        ("입력하면 실시간으로 필터됩니다", PT_DARK),
        ("", None),
        ("비워두면 전체 목록이 표시돼요", (120, 125, 130)),
    ]
    for i, (t, col) in enumerate(tips):
        if col:
            d.text((rx, 80 + i * 44), t, fill=col, font=pf(16, col == PT_DARK))

    # 화살표
    d.polygon([(px + pw + 10, H // 2 - 12),
                (px + pw + 36, H // 2),
                (px + pw + 10, H // 2 + 12)], fill=PT_TEAL)

    return img


def img_staging() -> Image.Image:
    """선택한 순서 패널."""
    W, H = 1100, 560
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    px, py, pw, ph = 40, 40, 500, 490
    d.rounded_rectangle([px, py, px + pw, py + ph], radius=12, fill=(255, 255, 255),
                          outline=(215, 218, 224))
    d.text((px + 16, py + 14), "선택한 순서  3개", fill=PT_DARK, font=pf(16, True))
    d.line([(px + 16, py + 40), (px + pw - 16, py + 40)], fill=(220, 222, 228))

    rows = [
        ("1", "주님 한 분만으로", False, True),
        ("2", "주의 이름 높이세", False, False),
        ("3", "로마서 8:28", True, False),
    ]
    for i, (num, title, bi, sel) in enumerate(rows):
        y = py + 50 + i * 130
        if sel:
            d.rectangle([px + 1, y, px + pw - 1, y + 120], fill=(237, 247, 255))
        d.text((px + 16, y + 40), num, fill=(60, 65, 70), font=pf(15))
        if bi:
            d.rounded_rectangle([px + 46, y + 34, px + 92, y + 58], radius=4,
                                  fill=(187, 222, 251))
            d.text((px + 52, y + 38), "성경", fill=(21, 100, 190), font=pf(12))
            d.text((px + 100, y + 36), title, fill=(30, 35, 40), font=pf(15))
        else:
            d.text((px + 46, y + 36), title, fill=(30, 35, 40), font=pf(15))
        # 버튼
        d.rounded_rectangle([px + pw - 90, y + 34, px + pw - 46, y + 60],
                              radius=5, fill=(240, 62, 62))
        d.text((px + pw - 80, y + 39), "✕ 제거", fill=PT_WHITE, font=pf(11))
        d.rounded_rectangle([px + pw - 40, y + 34, px + pw - 10, y + 60],
                              radius=5, fill=(208, 215, 220))
        d.text((px + pw - 36, y + 39), "⠿", fill=(80, 90, 95), font=pf(13))
        d.line([(px + 2, y + 120), (px + pw - 2, y + 120)], fill=(232, 235, 240))

    # 오른쪽 설명
    rx = 580
    items = [
        ("⠿  드래그로 순서 변경", "핸들을 위아래로 드래그"),
        ("✕  제거 버튼", "순서에서 삭제 (검색 목록 체크 해제)"),
        ("클릭", "오른쪽 미리보기가 해당 곡으로 전환"),
        ("패널 접기", "헤더 접기 아이콘으로 공간 확보"),
    ]
    for i, (action, desc) in enumerate(items):
        y = 80 + i * 90
        d.rounded_rectangle([rx, y, rx + 480, y + 74], radius=10, fill=(255, 255, 255),
                              outline=(215, 218, 224))
        d.text((rx + 16, y + 10), action, fill=PT_DARK, font=pf(14, True))
        d.text((rx + 16, y + 38), desc, fill=(100, 105, 112), font=pf(13))

    return img


def img_design() -> Image.Image:
    """디자인 패널 설명."""
    W, H = 1200, 600
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    # 슬라이드 미리보기 (왼쪽)
    mx, my, mw, mh = 40, 40, 460, 300
    d.rounded_rectangle([mx, my, mx + mw, my + mh], radius=8, fill=(27, 27, 27))
    d.text((mx + 80, my + 90), "주님 한 분만으로", fill=PT_WHITE, font=pf(24, True))
    d.text((mx + 80, my + 132), "Only You, Lord", fill=(255, 241, 118), font=pf(18))
    d.text((mx + 16, my + mh + 8), "▲ 실시간 미리보기 — 변경사항이 바로 반영돼요",
           fill=PT_TEAL, font=pf(12))

    # 옵션 목록 (왼쪽 하단)
    opts = [
        ("배경색", "■  #1B1B1B  (어두운 검정)"),
        ("글자 크기", "30 pt  [−] [+] 로 조절"),
        ("글자색", "■  흰색 / 노란색 등 선택"),
        ("영어 가사", "켜기/끄기 (토글)"),
        ("글자 위치", "상단 · 중단 · 하단"),
        ("글자 정렬", "좌측 · 중앙 · 우측"),
    ]
    for i, (lbl, val) in enumerate(opts):
        y = 380 + i * 34
        d.rounded_rectangle([mx, y, mx + mw, y + 28], radius=5,
                              fill=(255, 255, 255), outline=(220, 222, 228))
        d.text((mx + 12, y + 6), lbl, fill=PT_DARK, font=pf(13, True))
        d.text((mx + 130, y + 6), val, fill=(55, 60, 68), font=pf(13))

    # 오른쪽: 색상 스와치 예시
    rx = 540
    d.text((rx, 40), "배경색 팔레트", fill=PT_DARK, font=pf(16, True))
    swatches = [
        (0x1B, 0x1B, 0x1B), (0x12, 0x12, 0x12), (0x0F, 0x4C, 0x5C),
        (0x0B, 0x13, 0x2B), (0x5F, 0x0F, 0x40), (0xF4, 0xF1, 0xEA),
        (0xFF, 0xFF, 0xFF), (0x44, 0x5D, 0x48),
    ]
    for i, (r, g, b) in enumerate(swatches):
        x = rx + (i % 4) * 90
        y = 76 + (i // 4) * 90
        d.ellipse([x, y, x + 68, y + 68], fill=(r, g, b),
                   outline=(180, 185, 192))
        if (r, g, b) == (0x1B, 0x1B, 0x1B):
            d.text((x + 14, y + 22), "기본", fill=PT_WHITE, font=pf(13))

    d.text((rx, 278), "글자색 팔레트", fill=PT_DARK, font=pf(16, True))
    text_sw = [
        (0xFF, 0xFF, 0xFF), (0x00, 0x00, 0x00), (0xFF, 0xF1, 0x76),
        (0xFF, 0xCD, 0xD2), (0xC8, 0xE6, 0xC9), (0xD1, 0xC4, 0xE9),
    ]
    for i, (r, g, b) in enumerate(text_sw):
        x = rx + (i % 4) * 90
        y = 314 + (i // 4) * 90
        d.ellipse([x, y, x + 68, y + 68], fill=(r, g, b),
                   outline=(180, 185, 192))

    # PPTX 저장 버튼
    bx, by = rx, 510
    d.rounded_rectangle([bx, by, bx + 580, by + 64], radius=10, fill=PT_TEAL)
    d.text((bx + 200, by + 16), "PPTX 저장", fill=PT_WHITE, font=pf(22, True))
    d.text((rx, by + 82), "→  원하는 위치에 파일명을 입력하고 저장하면 완료!",
           fill=(85, 90, 98), font=pf(13))

    return img


def img_present_mode() -> Image.Image:
    """발표 모드 설명."""
    W, H = 1200, 560
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    # 컨트롤 바
    d.rounded_rectangle([30, 30, W - 30, 96], radius=10, fill=(196, 230, 234))
    d.text((50, 48), "■  발표 중  |  주님 한 분만으로", fill=PT_DARK, font=pf(17, True))
    d.text((680, 48), "◀", fill=PT_DARK, font=pf(18, True))
    d.rounded_rectangle([720, 42, 840, 82], radius=6, fill=PT_WHITE)
    d.text((748, 50), "3 / 12", fill=PT_DARK, font=pf(17, True))
    d.text((860, 48), "▶", fill=PT_DARK, font=pf(18, True))
    d.rounded_rectangle([920, 42, W - 38, 82], radius=6, fill=(228, 80, 80))
    d.text((950, 50), "■  발표 종료", fill=PT_WHITE, font=pf(15, True))

    # 슬라이드 목록
    names = ["1. 주님\n한분만으로", "2. 주님\n한분만으로", "3. 주님\n한분만으로",
             "4. 주의 이름\n높이세", "5. 로마서\n8:28"]
    for i, name in enumerate(names):
        x = 30 + i * 224
        col = PT_TEAL if i == 2 else (204, 212, 218)
        tc = PT_WHITE if i == 2 else (72, 82, 90)
        d.rounded_rectangle([x, 108, x + 212, 230], radius=8, fill=col)
        for j, line in enumerate(name.split('\n')):
            d.text((x + 14, 140 + j * 26), line, fill=tc, font=pf(14))
        if i == 2:
            d.rounded_rectangle([x + 6, 212, x + 56, 228], radius=3,
                                  fill=(255, 255, 255, 80))
            d.text((x + 10, 214), "현재", fill=PT_WHITE, font=pf(10))

    # 키보드 안내
    keys = [
        ("←  /  →", "이전 / 다음 슬라이드"),
        ("↑  /  ↓", "이전 / 다음 슬라이드"),
        ("숫자 + Enter", "해당 번호 슬라이드로 바로 이동"),
        ("ESC", "발표 종료"),
    ]
    kx = 30
    d.text((kx, 252), "키보드 단축키", fill=PT_DARK, font=pf(16, True))
    for i, (k, v) in enumerate(keys):
        y = 286 + i * 54
        d.rounded_rectangle([kx, y, kx + 560, y + 44], radius=8,
                              fill=(255, 255, 255), outline=(215, 218, 224))
        d.rounded_rectangle([kx + 8, y + 8, kx + 190, y + 36], radius=5,
                              fill=PT_DARK)
        d.text((kx + 20, y + 11), k, fill=PT_WHITE, font=pf(13, True))
        d.text((kx + 204, y + 12), v, fill=(48, 52, 60), font=pf(13))

    # 오른쪽: 슬라이드 편집 안내
    ex = 630
    d.text((ex, 252), "발표 중 슬라이드 즉석 편집", fill=PT_DARK, font=pf(16, True))
    edit_steps = [
        ("1", "슬라이드 목록에서\n수정할 슬라이드 더블클릭"),
        ("2", "가사 내용 수정"),
        ("3", "저장 클릭\n→ 발표 화면 즉시 반영"),
    ]
    for i, (num, txt) in enumerate(edit_steps):
        y = 294 + i * 84
        d.ellipse([ex, y, ex + 36, y + 36], fill=PT_TEAL)
        d.text((ex + 10, y + 6), num, fill=PT_WHITE, font=pf(14, True))
        for j, line in enumerate(txt.split('\n')):
            d.text((ex + 48, y + 4 + j * 22), line, fill=(48, 52, 60), font=pf(13))

    d.rounded_rectangle([ex, 548 - 30, ex + 540, 548], radius=6,
                          fill=(232, 248, 248), outline=PT_TEAL)
    d.text((ex + 14, 528), "📌  편집 내용은 현재 세션에만 적용돼요",
           fill=PT_DARK, font=pf(12))

    return img


def img_faq() -> Image.Image:
    """자주 묻는 질문."""
    W, H = 1200, 580
    img = Image.new('RGB', (W, H), (244, 247, 250))
    d = ImageDraw.Draw(img)

    faqs = [
        ("가져온 곡이 없다고 나와요",
         "선택한 폴더에 .pptx/.ppt 파일이 있는지 확인하세요"),
        ("중복 N개 건너뜀 메시지",
         "이미 저장된 곡이 포함된 경우 — 정상 동작이에요"),
        ("발표 화면이 주 모니터에 나와요",
         "OS 디스플레이 설정에서 '확장 모드'로 변경하세요"),
        ("영어 가사 빈 줄이 보여요",
         "디자인 패널 → '영어 가사 포함'을 꺼주세요"),
        ("곡을 전부 지우고 싶어요",
         "검색 패널 메뉴 → '전체 초기화' 버튼 클릭"),
        ("앱을 새로 받고 싶어요",
         "상단 헤더의 업데이트 배너 또는 ↺ 버튼을 확인하세요"),
    ]

    cols = 2
    bw = (W - 60) // cols - 20
    bh = 100

    for i, (q, a) in enumerate(faqs):
        col = i % cols
        row = i // cols
        x = 30 + col * (bw + 20)
        y = 20 + row * (bh + 20)
        d.rounded_rectangle([x, y, x + bw, y + bh], radius=10,
                              fill=(255, 255, 255), outline=(215, 218, 224))
        d.text((x + 14, y + 12), "Q.  " + q, fill=PT_DARK, font=pf(14, True))
        d.text((x + 14, y + 46), "→  " + a, fill=(65, 70, 78), font=pf(13))

    return img


# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 빌더
# ══════════════════════════════════════════════════════════════════════════════

def build_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    # ── 1. 표지 ───────────────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_DARK)
    # 그라디언트 사각형 (오른쪽에 밝은 틸 포인트)
    r2 = rect(s, SW * 6 // 10, Cm(0), SW * 4 // 10, SH, C_TEAL)
    r2.fill.fore_color.rgb = C_TEAL

    # 왼쪽 반투명 블록
    r3 = rect(s, Cm(0), Cm(0), SW * 6 // 10, SH, C_DARK)

    # 타이틀
    textbox(s, "예배 슬라이드",
            Cm(2), Cm(3), Cm(18), Cm(3.5),
            font_size=48, bold=True, color=C_WHITE)
    textbox(s, "사용자 가이드",
            Cm(2), Cm(6.8), Cm(18), Cm(2.5),
            font_size=36, bold=False, color=C_LTEAL)
    textbox(s, "찬양 · 성경 슬라이드를 쉽고 빠르게",
            Cm(2), Cm(10), Cm(22), Cm(1.5),
            font_size=16, color=RGBColor(0xA0, 0xD8, 0xD8))
    textbox(s, "v1.0.12",
            Cm(2), Cm(16.5), Cm(10), Cm(1.2),
            font_size=13, color=RGBColor(0x70, 0xA0, 0xA0))

    # ── 2. 프로그램 소개 ──────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "이런 걸 할 수 있어요")

    ov = img_app_overview()
    add_image(s, ov, Cm(1), Cm(3.5), SW - Cm(2))

    # ── 3. 화면 구성 ──────────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "화면 구성 한눈에 보기")

    ms = img_main_screen()
    add_image(s, ms, Cm(0.5), Cm(2.8), SW - Cm(1))

    # 레이블 텍스트박스
    textbox(s, "① 선택한 순서",
            Cm(2), Cm(16), Cm(8), Cm(1),
            font_size=11, bold=True, color=C_ACCENT)
    textbox(s, "② 찬양 검색 / 성경 탭",
            Cm(2), Cm(17), Cm(8), Cm(1),
            font_size=11, bold=True, color=C_ACCENT)
    textbox(s, "③ PPTX 디자인 & 저장",
            Cm(22), Cm(16), Cm(10), Cm(1),
            font_size=11, bold=True, color=C_ACCENT)

    # ── 4. 찬양 가져오기 ──────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "찬양 가져오기 — PPT 폴더 등록")

    ii = img_import_steps()
    add_image(s, ii, Cm(1), Cm(3.0), SW - Cm(2))

    textbox(s, "한 번 등록한 곡은 앱을 껐다 켜도 그대로 남아있어요",
            Cm(2), Cm(17.2), SW - Cm(4), Cm(1),
            font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── 5. 찬양 검색 & 선택 ──────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "찬양 검색 & 선택")

    ss = img_search_select()
    add_image(s, ss, Cm(1), Cm(3.0), SW - Cm(2))

    # ── 6. 선택한 순서 관리 ───────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "선택한 순서 관리")

    stg = img_staging()
    add_image(s, stg, Cm(1), Cm(3.0), SW - Cm(2))

    # ── 7. 디자인 & 저장 ─────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "디자인 설정 & PPTX 저장")

    di = img_design()
    add_image(s, di, Cm(1), Cm(3.0), SW - Cm(2))

    # ── 8. 발표 모드 ─────────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "화면 발표 모드")

    pm = img_present_mode()
    add_image(s, pm, Cm(1), Cm(3.0), SW - Cm(2))

    # ── 9. 성경 본문 추가 ────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "성경 본문 추가")
    _bible_slide_content(s)

    # ── 10. FAQ ──────────────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_LGRAY)
    _slide_header(s, "자주 묻는 질문")

    fq = img_faq()
    add_image(s, fq, Cm(1), Cm(3.0), SW - Cm(2))

    # ── 11. 마무리 ────────────────────────────────────────────────────────────
    s = new_slide(prs)
    bg(s, C_DARK)
    r2 = rect(s, SW * 6 // 10, Cm(0), SW * 4 // 10, SH, C_TEAL)

    textbox(s, "준비 완료!",
            Cm(2), Cm(5), Cm(20), Cm(3),
            font_size=44, bold=True, color=C_WHITE)
    textbox(s, "예배를 풍성하게 준비하세요",
            Cm(2), Cm(9.5), Cm(22), Cm(2),
            font_size=22, color=C_LTEAL)
    textbox(s, "문의 및 최신 버전  →  github.com/yjchae/make_ppt",
            Cm(2), Cm(16.5), Cm(26), Cm(1),
            font_size=12, color=RGBColor(0x80, 0xA8, 0xA8))

    return prs


def _slide_header(slide, title: str):
    """슬라이드 상단 헤더 바."""
    r = rect(slide, Cm(0), Cm(0), SW, Cm(2.6), C_DARK)
    textbox(slide, title,
            Cm(1.2), Cm(0.2), SW - Cm(2), Cm(2.2),
            font_size=22, bold=True, color=C_WHITE)
    # 하단 틸 라인
    rect(slide, Cm(0), Cm(2.6), SW, Cm(0.12), C_TEAL)


def _bible_slide_content(slide):
    """성경 본문 추가 슬라이드 — 텍스트 카드로 구성."""
    steps = [
        ("1단계", "성경 데이터 불러오기\n(최초 1회)", "상단 [성경 불러오기] 버튼 클릭\n→ JSON 파일 선택 → 버전명 입력"),
        ("2단계", "성경 본문 탭 클릭", "검색 패널 상단 탭에서\n[성경 본문] 탭 선택"),
        ("3단계", "책 · 장 선택", "드롭다운에서 책과 장을 선택하면\n절 목록이 표시됩니다"),
        ("4단계", "[+ 추가] 클릭", "원하는 절 오른쪽의 [+ 추가] 버튼을\n누르면 선택 순서에 추가됩니다"),
    ]
    bw = Cm(7.5)
    bh = Cm(12)
    gap = Cm(0.8)
    total = len(steps) * bw + (len(steps) - 1) * gap
    sx = (SW - total) / 2

    colors = [C_DARK, C_ACCENT, C_TEAL, RGBColor(0x12, 0x96, 0x74)]
    for i, ((step, title, desc), color) in enumerate(zip(steps, colors)):
        x = sx + i * (bw + gap)
        y = Cm(3.0)
        r = rect(slide, x, y, bw, bh, color)
        # 단계 번호
        pill(slide, step, x + Cm(0.3), y + Cm(0.3), Cm(2), Cm(0.7),
             RGBColor(0xFF, 0xFF, 0xFF),
             text_color=color, font_size=11, bold=True)
        textbox(slide, title,
                x + Cm(0.3), y + Cm(1.2), bw - Cm(0.6), Cm(2.5),
                font_size=15, bold=True, color=C_WHITE)
        textbox(slide, desc,
                x + Cm(0.3), y + Cm(3.5), bw - Cm(0.6), Cm(5),
                font_size=13, color=C_LTEAL, wrap=True)


# ── 메인 ─────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    prs = build_prs()
    out = '/home/user/make_ppt/예배슬라이드_사용자가이드.pptx'
    prs.save(out)
    print(f'저장 완료: {out}')
